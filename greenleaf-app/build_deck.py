# -*- coding: utf-8 -*-
"""
Builds GreenLeaf_Precision_Agriculture_Deck.pptx
Mirrors the consulting style of Hackathon-Slide-Deck.pdf (section nav tabs, action
titles, accent bar, callout boxes, charts, key-takeaway banners) but with our
GreenLeaf precision-agriculture data and topic-relevant imagery.
"""
import json, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn
from PIL import Image, ImageDraw, ImageEnhance

# ---------- paths ----------
ROOT = os.path.dirname(os.path.abspath(__file__))
ASSETS = r"C:\Users\navte\.cursor\projects\c-Users-navte-Downloads-greenleaf-app\assets"
DATA = json.load(open(os.path.join(ROOT, "greenleaf-app", "src", "data.json"), encoding="utf-8"))
OUT = os.path.join(ROOT, "GreenLeaf_Precision_Agriculture_Deck.pptx")

# ---------- palette ----------
INK        = RGBColor(0x1A, 0x1A, 0x1A)
INK_SOFT   = RGBColor(0x70, 0x6A, 0x60)
GREEN      = RGBColor(0x1C, 0x4A, 0x32)   # deep primary
GREEN_BR   = RGBColor(0x2D, 0x6A, 0x47)   # bright accent / nav active
GOLD       = RGBColor(0xC8, 0x99, 0x33)
TERRA      = RGBColor(0xB8, 0x42, 0x1B)
DANGER     = RGBColor(0xA8, 0x32, 0x32)
CREAM      = RGBColor(0xF4, 0xED, 0xE0)
GRAYBOX    = RGBColor(0xF1, 0xEC, 0xE2)
LINE       = RGBColor(0xD8, 0xCF, 0xBE)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BANNER     = GREEN

# ---------- derived numbers ----------
m = DATA["meta"]
plots = DATA["plots"]
prec_cost_total = sum(p["precision_cost"] for p in plots)
per_dollar = m["total_precision_benefit"] / prec_cost_total if prec_cost_total else 0
prec_share = m["total_precision_benefit"] / m["total_profit"]
treat = {t["treatment"]: t for t in DATA["treatments"]}
hl = treat["High Light"]; ctrl = treat["Control"]; hn = treat["High N"]
buckets = DATA["precision_buckets"]
top_tier = max(buckets, key=lambda b: b["avg_profit"])
base_tier = buckets[0]

# ---------- image preprocessing (circular crops + title scrim) ----------
def circle_crop(name, size=420, ring=14):
    src = os.path.join(ASSETS, name); dst = os.path.join(ASSETS, "c_" + name)
    im = Image.open(src).convert("RGB")
    w, h = im.size; s = min(w, h)
    im = im.crop(((w - s) // 2, (h - s) // 2, (w - s) // 2 + s, (h - s) // 2 + s)).resize((size, size))
    mask = Image.new("L", (size, size), 0)
    ImageDraw.Draw(mask).ellipse((0, 0, size, size), fill=255)
    out = Image.new("RGBA", (size, size), (0, 0, 0, 0)); out.paste(im, (0, 0), mask)
    ImageDraw.Draw(out).ellipse((ring // 2, ring // 2, size - ring // 2, size - ring // 2),
                                outline=(255, 255, 255), width=ring)
    out.save(dst); return dst

def title_bg():
    src = os.path.join(ASSETS, "hero_greenhouse.png"); dst = os.path.join(ASSETS, "hero_title.png")
    im = ImageEnhance.Brightness(Image.open(src).convert("RGB")).enhance(0.92).convert("RGBA")
    W, H = im.size; ov = Image.new("RGBA", (W, H), (0, 0, 0, 0)); d = ImageDraw.Draw(ov)
    for y in range(H):
        a = int(165 * max(0.0, (y - H * 0.5) / (H * 0.5)))
        d.line([(0, y), (W, y)], fill=(12, 22, 16, a))
    for x in range(int(W * 0.5)):
        a = int(120 * (1 - x / (W * 0.5)))
        d.line([(x, 0), (x, H)], fill=(12, 22, 16, a))
    Image.alpha_composite(im, ov).convert("RGB").save(dst); return dst

def darken(name, factor=0.78):
    src = os.path.join(ASSETS, name); dst = os.path.join(ASSETS, "d_" + name)
    ImageEnhance.Brightness(Image.open(src).convert("RGB")).enhance(factor).save(dst); return dst

HERO_T = title_bg()
BANNER_D = darken("greenhouse_banner.png", 0.72)
AERIAL_D = darken("farm_aerial.png", 0.74)
CIRC = {c: circle_crop(f"crop_{c}.png") for c in ("strawberry", "tomato", "pepper")}
LOGO = os.path.join(ASSETS, "greenleaf_logo.png")

# ---------- presentation ----------
prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def slide():
    return prs.slides.add_slide(BLANK)

def bg(s, color):
    s.background.fill.solid(); s.background.fill.fore_color.rgb = color

def _set_font(run, size, color, bold=False, italic=False, font="Calibri"):
    run.font.size = Pt(size); run.font.color.rgb = color; run.font.bold = bold
    run.font.italic = italic; run.font.name = font

def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=2, line_spacing=1.0, wrap=True):
    """runs: list of paragraphs, each a list of (txt, size, color, bold, italic) tuples."""
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after); p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for seg in para:
            txt, size, color, bold, italic = (seg + (False, False))[:5]
            r = p.add_run(); r.text = txt; _set_font(r, size, color, bold, italic)
    return tb

def rect(s, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE,
         dash=None, radius=None):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
        if dash:
            ln = sp.line._get_or_add_ln()
            d = ln.find(qn('a:prstDash'))
            if d is None:
                d = ln.makeelement(qn('a:prstDash'), {}); ln.append(d)
            d.set('val', dash)
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try: sp.adjustments[0] = radius
        except Exception: pass
    return sp

def shape_text(sp, runs, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, wrap=True):
    tf = sp.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left = Pt(6); tf.margin_right = Pt(6); tf.margin_top = Pt(3); tf.margin_bottom = Pt(3)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(1); p.space_before = Pt(0); p.line_spacing = 1.0
        for seg in para:
            txt, size, color, bold, italic = (seg + (False, False))[:5]
            r = p.add_run(); r.text = txt; _set_font(r, size, color, bold, italic)

def pic(s, path, x, y, w=None, h=None):
    return s.shapes.add_picture(path, x, y, width=w, height=h)

# ----- shared chrome -----
def header(s, eyebrow, title, title_size=21):
    rect(s, Inches(0.0), Inches(0.18), Inches(0.16), Inches(0.78), fill=GREEN_BR)
    text(s, Inches(0.38), Inches(0.16), Inches(11.0), Inches(0.32),
         [[(eyebrow, 11, INK_SOFT, False)]])
    text(s, Inches(0.38), Inches(0.44), Inches(12.0), Inches(0.62),
         [[(title, title_size, INK, True)]], line_spacing=0.95)
    pic(s, LOGO, Inches(12.35), Inches(0.12), w=Inches(0.85))

def navbar(s, active, tabs=None):
    tabs = tabs or ["Overview", "Analysis", "Strategy", "Conclusion"]
    strip_w = Inches(0.95 + 1.55 * len(tabs)); strip_x = (SW - strip_w) // 2
    y = Inches(6.92)
    rect(s, strip_x, y, strip_w, Inches(0.42), fill=GRAYBOX, line=LINE, line_w=0.75)
    tw = Inches(1.55); x = strip_x + Inches(0.475)
    for t in tabs:
        is_a = (t == active)
        text(s, x, y + Inches(0.07), tw, Inches(0.28),
             [[(t, 11, GREEN_BR if is_a else INK, is_a)]], align=PP_ALIGN.CENTER)
        if is_a:
            rect(s, x + Inches(0.2), y + Inches(0.355), tw - Inches(0.4), Pt(3), fill=GREEN_BR)
        x += tw
    text(s, SW - Inches(0.7), Inches(7.05), Inches(0.5), Inches(0.3),
         [[(str(s_num[0]), 11, INK_SOFT)]], align=PP_ALIGN.RIGHT)

s_num = [0]
def newslide(bgcolor=WHITE):
    s_num[0] += 1; s = slide(); bg(s, bgcolor); return s

# ----- chart styling -----
def style_axes(chart, val_fmt=None, cat_size=10, val_size=9, val_axis=True):
    try: chart.has_title = False
    except Exception: pass
    try:
        ca = chart.category_axis
        ca.tick_labels.font.size = Pt(cat_size); ca.tick_labels.font.color.rgb = INK
        ca.format.line.color.rgb = LINE
        ca.has_major_gridlines = False; ca.has_minor_gridlines = False
    except Exception: pass
    try:
        va = chart.value_axis
        va.visible = val_axis
        va.has_major_gridlines = True
        va.major_gridlines.format.line.color.rgb = RGBColor(0xEC, 0xE6, 0xDA)
        va.major_gridlines.format.line.width = Pt(0.5)
        va.tick_labels.font.size = Pt(val_size); va.tick_labels.font.color.rgb = INK_SOFT
        va.format.line.fill.background()
        if val_fmt: va.tick_labels.number_format = val_fmt; va.tick_labels.number_format_is_linked = False
    except Exception: pass

def color_series(series, color):
    series.format.fill.solid(); series.format.fill.fore_color.rgb = color
    series.format.line.fill.background()

def labels(series, fmt=None, size=10, color=INK, bold=True, pos=XL_LABEL_POSITION.OUTSIDE_END):
    series.has_data_labels = True
    dl = series.data_labels
    dl.font.size = Pt(size); dl.font.color.rgb = color; dl.font.bold = bold
    if fmt: dl.number_format = fmt; dl.number_format_is_linked = False
    try: dl.position = pos
    except Exception: pass

def panel_title(s, x, y, w, label):
    sp = rect(s, x, y, w, Inches(0.34), fill=WHITE, line=GREEN_BR, line_w=1.25)
    shape_text(sp, [[(label, 12, GREEN, True)]])

# =====================================================================
# SLIDE 1 — TITLE
# =====================================================================
s = newslide()
pic(s, HERO_T, 0, 0, w=SW, h=SH)
# logo lockup in white box
box = rect(s, Inches(0.7), Inches(2.55), Inches(3.0), Inches(1.25), fill=WHITE, radius=0.08,
           shape=MSO_SHAPE.ROUNDED_RECTANGLE)
pic(s, LOGO, Inches(0.85), Inches(2.68), w=Inches(1.0))
text(s, Inches(1.85), Inches(2.78), Inches(1.8), Inches(1.0),
     [[("Green", 23, GREEN, True), ("Leaf", 23, GOLD, True)],
      [("PRECISION CEA", 9.5, INK_SOFT, True)]], anchor=MSO_ANCHOR.MIDDLE)
# vertical divider
rect(s, Inches(3.95), Inches(2.5), Pt(3), Inches(1.4), fill=WHITE)
# big title
text(s, Inches(4.25), Inches(2.5), Inches(8.6), Inches(1.7),
     [[("Precision Agriculture:", 36, WHITE, True)],
      [("Does the Spend Pay Off?", 36, WHITE, True)]], line_spacing=1.02)
# bottom-left meta
text(s, Inches(0.7), Inches(6.05), Inches(6.0), Inches(0.9),
     [[("Presenting to: ", 15, WHITE, True), ("RBC — Capital Allocation Review", 15, WHITE, False)],
      [("Date: ", 15, WHITE, True), ("May 30, 2026", 15, WHITE, False)]], space_after=4)
# bottom-right crop circles (replacing headshots) + labels
text(s, Inches(7.55), Inches(6.7), Inches(1.6), Inches(0.3),
     [[("Crops studied:", 12, WHITE, True)]], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
cx = Inches(9.3)
for c in ("strawberry", "tomato", "pepper"):
    pic(s, CIRC[c], cx, Inches(6.35), w=Inches(0.95), h=Inches(0.95))
    cx += Inches(1.25)
text(s, Inches(9.3), Inches(6.0), Inches(3.75), Inches(0.3),
     [[("Strawberry · Tomato · Pepper / Cucumber", 9.5, WHITE, False)]], align=PP_ALIGN.LEFT)

# =====================================================================
# helper for the Executive-Summary / Conclusion grid (slides 2 & 11)
# =====================================================================
def exec_summary(s, eyebrow):
    header(s, eyebrow, "Executive Summary")
    lx = Inches(0.4); lw = Inches(1.55); cx = Inches(2.1); cw = Inches(10.85)
    rows = [(1.15, "Objective"), (2.15, "Problems"), (3.45, "Strategy"), (5.35, "Key\nTakeaway")]
    for yy, lbl in rows:
        rect(s, cx - Inches(0.18), Inches(yy), Pt(2.5), Inches(0.7 if lbl != "Strategy" else 1.6),
             fill=INK)
        text(s, lx, Inches(yy), lw, Inches(0.9),
             [[(part, 14, INK, True)] for part in lbl.split("\n")],
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0)
    # Objective box
    ob = rect(s, cx, Inches(1.1), cw, Inches(0.78), fill=GRAYBOX)
    shape_text(ob, [[("Determine whether GreenLeaf's precision-agriculture spending delivers measurable "
                      "financial returns across 120 microplots — and where capital should flow next.", 13, INK, False)]])
    # Problems (3 green boxes)
    pbw = (cw - Inches(0.6)) / 3
    probs = ["6,460 alerts in-season — only 74.6% were acted on",
             "High-stress events spike in spring, peaking the week of Apr 5",
             "Spend is spread evenly across treatments, not by ROI"]
    px = cx
    for p in probs:
        b = rect(s, px, Inches(2.15), pbw, Inches(0.92), fill=WHITE, line=GREEN_BR, line_w=1.5)
        shape_text(b, [[(p, 12, INK, False)]])
        px += pbw + Inches(0.3)
    # Strategy banner
    bnx, bny, bnw, bnh = cx, Inches(3.45), cw, Inches(1.65)
    pic(s, BANNER_D, bnx, bny, w=bnw, h=bnh)
    hdr = rect(s, bnx, bny, bnw, Inches(0.42), fill=GREEN)
    shape_text(hdr, [[("Smart Precision", 15, WHITE, True)]])
    half = bnw / 2
    text(s, bnx, bny + Inches(0.55), half, Inches(1.0),
         [[("Returns", 16, WHITE, True)],
          [("66% of season profit traces to precision actions", 12, WHITE, False)]],
         align=PP_ALIGN.CENTER, space_after=4)
    text(s, bnx + half, bny + Inches(0.55), half, Inches(1.0),
         [[("Resilience", 16, WHITE, True)],
          [("Targeted, faster action cuts crop-stress exposure", 12, WHITE, False)]],
         align=PP_ALIGN.CENTER, space_after=4)
    # Key takeaway
    kt = rect(s, cx, Inches(5.35), cw, Inches(1.15), fill=GRAYBOX)
    shape_text(kt, [[("Precision spending is worth financing: ", 13, INK, True),
                     ("$47.0K of $71.7K season profit (66%%) is attributable to precision actions — "
                      "about $%.2f back per $1 spent. Reallocate toward High-Light blocks and integrated "
                      "pest management." % per_dollar, 13, INK, False)]])

# =====================================================================
# SLIDE 2 — EXECUTIVE SUMMARY (Overview)
# =====================================================================
s = newslide(); exec_summary(s, "Overview"); navbar(s, "Overview")

# =====================================================================
# SLIDE 3 — METHODOLOGY (Analysis)
# =====================================================================
s = newslide()
header(s, "Analysis", "How we tested the precision-agriculture investment case")
text(s, Inches(0.9), Inches(1.35), Inches(4.0), Inches(0.4), [[("Methodology", 16, INK, True)]])
steps = [
    ("Season Economics", "Per-plot revenue, cost, profit and ROI from season_summary", GREEN),
    ("Stress & Alerts", "Daily plant-stress index, alert flags and action response", GREEN_BR),
    ("Precision Spend", "Share of plot cost spent on alert-triggered (precision) actions", GOLD),
    ("Treatment ROI", "120 plots across 7 treatments vs the Control baseline", TERRA),
]
fy = Inches(1.95); fh = Inches(0.72); maxw = Inches(4.6); minw = Inches(2.6)
for i, (t, _, col) in enumerate(steps):
    w = Emu(int(maxw) - int((int(maxw) - int(minw)) * i / 3))
    x = Inches(0.9) + (maxw - w) / 2
    b = rect(s, x, fy, w, fh, fill=col)
    shape_text(b, [[(t, 13.5, WHITE, True)]])
    # connector + description
    rect(s, Inches(5.7), fy + fh / 2 - Pt(1), Inches(0.8), Pt(2), fill=col)
    text(s, Inches(6.7), fy, Inches(6.3), fh,
         [[(steps[i][1], 13, INK, False)]], anchor=MSO_ANCHOR.MIDDLE)
    fy += fh + Inches(0.34)
# bottom funnel tip
tip = rect(s, Inches(2.55) + (maxw)/2 - Inches(0.45), fy, Inches(0.9), Inches(0.5),
           fill=INK_SOFT, shape=MSO_SHAPE.ISOSCELES_TRIANGLE)
tip.rotation = 180
navbar(s, "Analysis")

# =====================================================================
# SLIDE 4 — SEASON AT A GLANCE (Overview/Analysis)
# =====================================================================
s = newslide()
header(s, "Analysis · Season at a Glance",
       "Eight farms turned $238K revenue into $71.7K profit at a 42.8% average ROI")
# KPI strip
kpis = [("$238K", "Season revenue", GREEN), ("$71.7K", "Season profit", GREEN_BR),
        ("$47.0K", "Precision benefit", GOLD), ("42.8%", "Avg plot ROI", TERRA),
        ("6,460", "Alerts → 74.6% acted", INK_SOFT)]
kx = Inches(0.4); kw = Inches(2.45)
for val, lbl, col in kpis:
    b = rect(s, kx, Inches(1.35), kw, Inches(1.0), fill=GRAYBOX, line=col, line_w=1.5, radius=0.06,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    shape_text(b, [[(val, 22, col, True)], [(lbl, 10.5, INK, False)]])
    kx += kw + Inches(0.13)
# farm chart
panel_title(s, Inches(0.4), Inches(2.7), Inches(7.4), "Profit & precision benefit by farm")
farms = sorted(DATA["farms"], key=lambda f: -f["total_profit"])
cd = CategoryChartData()
cd.categories = [f["farm_name"].replace("BC Harvest", "BC") for f in farms]
cd.add_series("Total profit", [round(f["total_profit"]) for f in farms])
cd.add_series("Precision benefit", [round(f["total_precision_benefit"]) for f in farms])
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.4), Inches(3.1),
                        Inches(7.4), Inches(3.55), cd).chart
gf.has_legend = True; gf.legend.position = XL_LEGEND_POSITION.BOTTOM
gf.legend.include_in_layout = False; gf.legend.font.size = Pt(9)
color_series(gf.series[0], GREEN); color_series(gf.series[1], GOLD)
style_axes(gf, val_fmt='"$"#,##0')
gf.value_axis.tick_labels.number_format = '"$"#,##0'; gf.value_axis.tick_labels.number_format_is_linked = False
# right callouts
text(s, Inches(8.1), Inches(2.85), Inches(4.9), Inches(0.4), [[("What stands out", 14, INK, True)]])
cinfo = [("BC Harvest 8 leads", "$21.2K profit · $14.0K precision benefit across 34 plots", GREEN),
         ("Strawberries pay best", "$679 avg profit/plot — highest of four crops", GOLD),
         ("Pest pressure dominates alerts", "99% of 6,460 alerts are High Pest Pressure", TERRA)]
cy = Inches(3.3)
for t, d, col in cinfo:
    b = rect(s, Inches(8.1), cy, Inches(4.9), Inches(0.95), fill=WHITE, line=col, line_w=1.25,
             dash="dash")
    shape_text(b, [[(t, 12.5, col, True)], [(d, 11, INK, False)]], align=PP_ALIGN.LEFT)
    cy += Inches(1.08)
navbar(s, "Analysis")

# =====================================================================
# SLIDE 5 — WEEKLY TRIAGE (Part A)
# =====================================================================
s = newslide()
header(s, "Analysis · Weekly Triage",
       "High-stress events cluster in spring — the week of Apr 5 is the moment to act")
panel_title(s, Inches(0.4), Inches(1.4), Inches(8.3), "High-stress plot-days per week (Feb 15 → Sep 12)")
wk = DATA["weekly"]
cd = CategoryChartData()
cd.categories = [w["week_label"][5:] for w in wk]
cd.add_series("High-stress events", [w["high_stress_events"] for w in wk])
ch = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.4), Inches(1.8),
                        Inches(8.3), Inches(4.7), cd).chart
ch.has_legend = False
sr = ch.series[0]; color_series(sr, GREEN_BR)
peak_i = max(range(len(wk)), key=lambda i: wk[i]["high_stress_events"])
sr.points[peak_i].format.fill.solid(); sr.points[peak_i].format.fill.fore_color.rgb = TERRA
style_axes(ch, cat_size=7)
# callouts
text(s, Inches(9.0), Inches(1.6), Inches(4.0), Inches(0.4), [[("Triage signal", 14, INK, True)]])
metrics = [("911", "high-stress plot-days all season", TERRA),
           ("74.6%", "of alerts converted to action", GREEN),
           ("0.75 d", "average action delay", GOLD),
           ("2,067", "same-day responses to alerts", GREEN_BR)]
my = Inches(2.05)
for v, l, col in metrics:
    b = rect(s, Inches(9.0), my, Inches(4.0), Inches(0.92), fill=GRAYBOX, line=col, line_w=1.5)
    shape_text(b, [[(v, 19, col, True)], [(l, 11, INK, False)]], align=PP_ALIGN.LEFT)
    my += Inches(1.05)
navbar(s, "Analysis")

# =====================================================================
# SLIDE 6 — PRECISION ROI HEADLINE (Part B)
# =====================================================================
s = newslide()
header(s, "Analysis · Precision ROI",
       "Two-thirds of season profit traces directly to precision actions")
# big stat block
big = rect(s, Inches(0.4), Inches(1.5), Inches(4.6), Inches(4.7), fill=GREEN)
shape_text(big, [[("66%", 78, WHITE, True)],
                 [("of total profit attributable", 14, WHITE, False)],
                 [("to precision actions", 14, WHITE, False)]])
# narrative + pair stats
text(s, Inches(5.3), Inches(1.55), Inches(7.6), Inches(1.6),
     [[("Of ", 17, INK, False), ("$71.7K", 17, GREEN, True),
       (" in season profit, ", 17, INK, False), ("$47.0K", 17, GOLD, True),
       (" traces directly to precision actions — about ", 17, INK, False),
       ("$%.2f" % per_dollar, 17, TERRA, True), (" back per $1 of precision spend." , 17, INK, False)]],
     line_spacing=1.15)
pairs = [("$47.0K", "precision benefit", GOLD), ("42.8%", "avg plot ROI", GREEN_BR),
         ("$%.2f" % per_dollar, "return per $1 spent", TERRA)]
pxx = Inches(5.3)
for v, l, col in pairs:
    b = rect(s, pxx, Inches(3.5), Inches(2.4), Inches(1.15), fill=GRAYBOX, line=col, line_w=1.5,
             radius=0.06, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    shape_text(b, [[(v, 24, col, True)], [(l, 11, INK, False)]])
    pxx += Inches(2.55)
fin = rect(s, Inches(5.3), Inches(4.95), Inches(7.5), Inches(1.25), fill=WHITE, line=GREEN_BR,
           line_w=1.5, dash="dash")
shape_text(fin, [[("Financing case: ", 13.5, GREEN, True),
                  ("precision agriculture is not a cost centre — it is the single largest identifiable "
                   "driver of GreenLeaf's profitability, and it scales with disciplined, alert-driven spending.",
                   13, INK, False)]], align=PP_ALIGN.LEFT)
navbar(s, "Analysis")

# =====================================================================
# SLIDE 7 — PRECISION TIERS (Part B)
# =====================================================================
s = newslide()
header(s, "Analysis · Precision Tiers",
       "The most precision-intensive plots earn the most — $692 vs $555 profit/plot")
panel_title(s, Inches(0.4), Inches(1.45), Inches(6.0), "Avg profit per plot by precision-spend tier")
cd = CategoryChartData(); cd.categories = [b["prec_bucket"] for b in buckets]
cd.add_series("Avg profit", [round(b["avg_profit"]) for b in buckets])
c1 = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.4), Inches(1.85),
                        Inches(6.0), Inches(4.5), cd).chart
c1.has_legend = False; sr = c1.series[0]
for i, b in enumerate(buckets):
    sr.points[i].format.fill.solid()
    sr.points[i].format.fill.fore_color.rgb = GOLD if b is top_tier else GREEN
labels(sr, fmt='"$"#,##0', size=10); style_axes(c1, val_fmt='"$"#,##0')
panel_title(s, Inches(6.7), Inches(1.45), Inches(6.0), "Avg ROI by precision-spend tier")
cd = CategoryChartData(); cd.categories = [b["prec_bucket"] for b in buckets]
cd.add_series("Avg ROI", [round(b["avg_roi"], 4) for b in buckets])
c2 = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(6.7), Inches(1.85),
                        Inches(6.0), Inches(4.5), cd).chart
c2.has_legend = False; sr2 = c2.series[0]
for i, b in enumerate(buckets):
    sr2.points[i].format.fill.solid()
    sr2.points[i].format.fill.fore_color.rgb = GOLD if b is top_tier else GREEN_BR
labels(sr2, fmt='0%', size=10); style_axes(c2, val_fmt='0%')
text(s, Inches(0.4), Inches(6.45), Inches(12.3), Inches(0.4),
     [[("Top tier (~20% of cost on precision): ", 12, GREEN, True),
       ("highest profit ($692/plot) AND highest ROI (48.5%) — more precision, more return.", 12, INK, False)]],
     align=PP_ALIGN.CENTER)
navbar(s, "Analysis")

# =====================================================================
# SLIDE 8 — TREATMENT AUDIT (Part C)
# =====================================================================
s = newslide()
header(s, "Analysis · Treatment Audit",
       "The 'more nitrogen = more profit' rule is a myth — High-Light wins on ROI")
panel_title(s, Inches(0.4), Inches(1.45), Inches(7.4), "Avg profit per plot by treatment")
tr = sorted(DATA["treatments"], key=lambda t: -t["avg_profit"])
cd = CategoryChartData(); cd.categories = [t["treatment"] for t in tr]
cd.add_series("Avg profit", [round(t["avg_profit"]) for t in tr])
c = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.4), Inches(1.85),
                       Inches(7.4), Inches(4.6), cd).chart
c.has_legend = False; sr = c.series[0]
for i, t in enumerate(tr):
    sr.points[i].format.fill.solid()
    col = GOLD if t["treatment"] == "High Light" else (GREEN if t["treatment"] == "High N" else GREEN_BR)
    sr.points[i].format.fill.fore_color.rgb = col
labels(sr, fmt='"$"#,##0', size=10, pos=XL_LABEL_POSITION.OUTSIDE_END)
style_axes(c, val_fmt='"$"#,##0')
# myth / truth
mt = rect(s, Inches(8.1), Inches(1.7), Inches(4.9), Inches(1.55), fill=GRAYBOX)
shape_text(mt, [[("COMMON WISDOM", 10.5, INK_SOFT, True)],
                [("\"Push the nitrogen — heavier feed means a heavier harvest and heavier cheques.\"",
                  12.5, INK, False, True)]], align=PP_ALIGN.LEFT)
tb = rect(s, Inches(8.1), Inches(3.4), Inches(4.9), Inches(1.85), fill=GREEN, )
shape_text(tb, [[("WHAT THE DATA SAYS", 10.5, RGBColor(0xCF,0xE0,0xD4), True)],
                [("High-N delivers the top yield (7.46 kg/m²) but only $651 profit. "
                  "High-Light earns $839 at 58% ROI on mid yield.", 12.5, WHITE, False)]],
         align=PP_ALIGN.LEFT)
rec = rect(s, Inches(8.1), Inches(5.4), Inches(4.9), Inches(0.95), fill=WHITE, line=GOLD, line_w=1.5,
           dash="dash")
shape_text(rec, [[("Yield ≠ profit. ", 12.5, GOLD, True),
                  ("Treat High-N as a yield play only, not a profit lever.", 12, INK, False)]],
           align=PP_ALIGN.LEFT)
navbar(s, "Analysis")

# =====================================================================
# SLIDE 9 — STRATEGY SHORT TERM
# =====================================================================
s = newslide()
header(s, "Strategy · Short-term",
       "Reallocate spend toward High-Light blocks and integrated pest; act same-day")
pic(s, BANNER_D, Inches(0.0), Inches(1.25), w=SW, h=Inches(5.4))
band = rect(s, Inches(0.0), Inches(3.05), SW, Inches(1.5), fill=GREEN)
shape_text(band, [[("Shift capital from negative-ROI treatments to High-Light and integrated pest, "
                    "and convert the 25% of pest alerts still missed into same-day action.", 17, WHITE, True)]])
calls = [("Expand High-Light", "+$258/plot vs Control · 58% ROute", GREEN_BR, Inches(0.6), Inches(1.55)),
         ("Scale back Shade & Reduced-Pest", "31% ROI · 14–19 high-stress days/plot", TERRA, Inches(0.6), Inches(5.0)),
         ("Act on pest alerts same-day", "75% acted today vs 0.75d avg delay", GOLD, Inches(8.9), Inches(5.0))]
for t, d, col, x, y in calls:
    b = rect(s, x, y, Inches(3.8), Inches(0.95), fill=WHITE, line=col, line_w=1.75, dash="dash")
    shape_text(b, [[(t, 13, col, True)], [(d.replace("ROute","ROI"), 11, INK, False)]], align=PP_ALIGN.LEFT)
navbar(s, "Strategy")

# =====================================================================
# SLIDE 10 — STRATEGY LONG TERM
# =====================================================================
s = newslide()
header(s, "Strategy · Long-term",
       "Scale precision financing across all eight farms and densify sensing")
pic(s, AERIAL_D, Inches(0.0), Inches(1.25), w=SW, h=Inches(5.4))
band = rect(s, Inches(0.0), Inches(3.15), SW, Inches(1.5), fill=GREEN)
shape_text(band, [[("Fund a fleet-wide precision program: replicate High-Light + integrated pest "
                    "across 120 plots and raise sensor density to close the alert-response gap.", 17, WHITE, True)]])
lt = [("Replicate the winners", "High-Light + Integrated Pest farm-wide", GREEN_BR, Inches(0.6), Inches(1.55)),
      ("Densify sensing", "Convert the 25% of unactioned alerts", GOLD, Inches(8.9), Inches(1.55)),
      ("Finance the spend", "~$%.2f returned per $1 of precision cost" % per_dollar, GREEN, Inches(8.9), Inches(5.05))]
for t, d, col, x, y in lt:
    b = rect(s, x, y, Inches(3.8), Inches(0.95), fill=WHITE, line=col, line_w=1.75, dash="dash")
    shape_text(b, [[(t, 13, col, True)], [(d, 11, INK, False)]], align=PP_ALIGN.LEFT)
navbar(s, "Strategy")

# =====================================================================
# SLIDE 11 — CONCLUSION (Executive Summary)
# =====================================================================
s = newslide(); exec_summary(s, "Conclusion"); navbar(s, "Conclusion")

# =====================================================================
# SLIDE 12 — APPENDIX DIVIDER
# =====================================================================
s = newslide(GREEN)
pic(s, AERIAL_D, 0, 0, w=SW, h=SH)
ov = rect(s, 0, 0, SW, SH, fill=GREEN); 
# fake transparency via lighter overlay not available; use band instead
ov.fill.background()
band = rect(s, 0, Inches(3.0), SW, Inches(1.5), fill=GREEN)
shape_text(band, [[("Appendix", 40, WHITE, True)]])
pic(s, LOGO, Inches(12.35), Inches(0.25), w=Inches(0.85))

# =====================================================================
# SLIDE 13 — APPENDIX: TREATMENT SCORECARD
# =====================================================================
s = newslide()
header(s, "Appendix", "Treatment scorecard — ranked by ROI vs Control baseline")
# left: key issue/result/recommendation
ki = rect(s, Inches(0.4), Inches(1.5), Inches(4.5), Inches(2.0), fill=WHITE, line=GREEN, line_w=2)
shape_text(ki, [[("KEY QUESTION", 11, GREEN, True)],
                [("Does precision spending earn its keep — and which treatments should we fund?",
                  12.5, INK, False)]], align=PP_ALIGN.LEFT)
res = rect(s, Inches(0.4), Inches(3.65), Inches(4.5), Inches(1.35), fill=GRAYBOX)
shape_text(res, [[("RESULT", 11, INK_SOFT, True)],
                 [("66% of profit from precision · High-Light best ROI · High-N is yield-only",
                   12, INK, False)]], align=PP_ALIGN.LEFT)
rc = rect(s, Inches(0.4), Inches(5.15), Inches(4.5), Inches(1.35), fill=GREEN)
shape_text(rc, [[("RECOMMENDATION", 11, RGBColor(0xCF,0xE0,0xD4), True)],
                [("Finance precision; expand High-Light & integrated pest; trim Shade & Reduced-Pest.",
                  12, WHITE, False)]], align=PP_ALIGN.LEFT)
# right: table
tr = sorted(DATA["treatments"], key=lambda t: -t["avg_roi"])
rows = len(tr) + 1
tbl = s.shapes.add_table(rows, 6, Inches(5.15), Inches(1.5), Inches(7.85), Inches(5.0)).table
hdrs = ["Treatment", "n", "Profit", "ROI", "Yield", "HS days"]
widths = [2.6, 0.7, 1.35, 1.1, 1.05, 1.05]
for j, wv in enumerate(widths): tbl.columns[j].width = Inches(wv)
for j, htxt in enumerate(hdrs):
    cell = tbl.cell(0, j); cell.fill.solid(); cell.fill.fore_color.rgb = GREEN
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.RIGHT
    r = p.add_run(); r.text = htxt; _set_font(r, 11.5, WHITE, True)
for i, t in enumerate(tr, start=1):
    vals = [t["treatment"], str(t["n_plots"]), "$%d" % round(t["avg_profit"]),
            "%.1f%%" % (t["avg_roi"] * 100), "%.2f" % t["avg_yield"], "%.1f" % t["avg_hs_days"]]
    best = t["treatment"] == "High Light"
    rowcol = RGBColor(0xF6,0xEE,0xD8) if best else (WHITE if i % 2 else GRAYBOX)
    for j, v in enumerate(vals):
        cell = tbl.cell(i, j); cell.fill.solid(); cell.fill.fore_color.rgb = rowcol
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.RIGHT
        r = p.add_run(); r.text = v
        _set_font(r, 11, GREEN if best and j == 0 else INK, best or j == 0)
navbar(s, "Conclusion", tabs=["Overview", "Analysis", "Strategy", "Conclusion", "Appendix"])

prs.save(OUT)
print("Saved:", OUT)
print("Slides:", len(prs.slides._sldIdLst))
print("per_dollar=%.3f prec_share=%.3f prec_cost_total=%.0f" % (per_dollar, prec_share, prec_cost_total))
